#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT/PDF 翻译工作台 - 后端 Flask 应用
=====================================
功能：
  1. 文件上传与解析（PPTX / PDF）
  2. 多 AI API 翻译（DeepSeek、豆包/字节豆包、OpenAI 兼容接口）
  3. 翻译记忆库（SQLite，支持相似度匹配）
  4. 语料库/术语管理
  5. 多格式导出（MD、JSON、DOCX、PDF）
  6. 简单的用户认证

依赖：
  pip install flask python-pptx PyPDF2 python-docx fpdf2

启动方式：
  python app.py
"""

import os
import sys
import json
import uuid
import time
import hashlib
import sqlite3
import base64
import re
import threading
import logging
import traceback
from datetime import datetime, timedelta
from functools import wraps
from io import BytesIO
from contextlib import contextmanager

from flask import (
    Flask, request, jsonify, session,
    send_file, send_from_directory, make_response,
    render_template
)

# ---------------------------------------------------------------------------
# 第三方库导入（可选，缺失时给出友好提示）
# ---------------------------------------------------------------------------
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from PyPDF2 import PdfReader
    HAS_PDF2 = True
except ImportError:
    HAS_PDF2 = False

try:
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from fpdf import FPDF
    HAS_FPDF = True
except ImportError:
    HAS_FPDF = False

# ---------------------------------------------------------------------------
# 日志配置
# ---------------------------------------------------------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
)
logger = logging.getLogger("ppt_translate")

# ---------------------------------------------------------------------------
# Flask 应用初始化
# ---------------------------------------------------------------------------
app = Flask(
    __name__,
    static_folder=None,
)
app.secret_key = os.environ.get("FLASK_SECRET_KEY", uuid.uuid4().hex)
app.config["MAX_CONTENT_LENGTH"] = 100 * 1024 * 1024  # 100 MB

# 上传文件存储目录（云环境使用 /tmp，本地使用项目目录）
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
if os.environ.get("RENDER") or os.environ.get("DYNO"):
    UPLOAD_FOLDER = "/tmp/uploads"
    DATABASE_PATH = "/tmp/translation.db"
else:
    UPLOAD_FOLDER = os.path.join(BASE_DIR, "uploads")
    DATABASE_PATH = os.path.join(BASE_DIR, "translation.db")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)


# ===========================================================================
# 数据库管理
# ===========================================================================

def get_db():
    """获取数据库连接（线程局部）。"""
    if not hasattr(threading.current_thread(), "_db_conn"):
        threading.current_thread()._db_conn = sqlite3.connect(DATABASE_PATH)
        threading.current_thread()._db_conn.row_factory = sqlite3.Row
        threading.current_thread()._db_conn.execute("PRAGMA journal_mode=WAL")
        threading.current_thread()._db_conn.execute("PRAGMA foreign_keys=ON")
    return threading.current_thread()._db_conn


@contextmanager
def db_cursor():
    """上下文管理器：获取游标，自动提交/回滚。"""
    conn = get_db()
    cursor = conn.cursor()
    try:
        yield cursor
        conn.commit()
    except Exception:
        conn.rollback()
        raise


def init_db():
    """初始化数据库表结构和默认数据。"""
    with db_cursor() as cur:
        # 用户表
        cur.execute("""
            CREATE TABLE IF NOT EXISTS users (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                username TEXT UNIQUE NOT NULL,
                password_hash TEXT NOT NULL,
                is_admin INTEGER DEFAULT 0,
                created_at TEXT DEFAULT (datetime('now','localtime'))
            )
        """)

        # 翻译记忆库
        cur.execute("""
            CREATE TABLE IF NOT EXISTS translation_memory (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER,
                source_lang TEXT NOT NULL,
                target_lang TEXT NOT NULL,
                source_text TEXT NOT NULL,
                target_text TEXT NOT NULL,
                similarity REAL DEFAULT 1.0,
                created_at TEXT DEFAULT (datetime('now','localtime')),
                updated_at TEXT DEFAULT (datetime('now','localtime')),
                FOREIGN KEY (user_id) REFERENCES users(id)
            )
        """)

        # 语料库/术语表
        cur.execute("""
            CREATE TABLE IF NOT EXISTS glossary (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER,
                source_term TEXT NOT NULL,
                target_term TEXT NOT NULL,
                source_lang TEXT NOT NULL,
                target_lang TEXT NOT NULL,
                category TEXT DEFAULT '',
                note TEXT DEFAULT '',
                created_at TEXT DEFAULT (datetime('now','localtime')),
                FOREIGN KEY (user_id) REFERENCES users(id)
            )
        """)

        # API 配置表
        cur.execute("""
            CREATE TABLE IF NOT EXISTS api_configs (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER,
                api_type TEXT NOT NULL,
                api_key TEXT DEFAULT '',
                api_url TEXT DEFAULT '',
                model TEXT DEFAULT '',
                is_default INTEGER DEFAULT 0,
                created_at TEXT DEFAULT (datetime('now','localtime')),
                FOREIGN KEY (user_id) REFERENCES users(id)
            )
        """)

        # 翻译历史表
        cur.execute("""
            CREATE TABLE IF NOT EXISTS translation_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER,
                source_text TEXT NOT NULL,
                translated_text TEXT NOT NULL,
                source_lang TEXT NOT NULL,
                target_lang TEXT NOT NULL,
                api_type TEXT DEFAULT '',
                model TEXT DEFAULT '',
                use_tm INTEGER DEFAULT 1,
                created_at TEXT DEFAULT (datetime('now','localtime')),
                FOREIGN KEY (user_id) REFERENCES users(id)
            )
        """)

        # 创建索引
        cur.execute("CREATE INDEX IF NOT EXISTS idx_tm_source ON translation_memory(source_text)")
        cur.execute("CREATE INDEX IF NOT EXISTS idx_tm_langs ON translation_memory(source_lang, target_lang)")
        cur.execute("CREATE INDEX IF NOT EXISTS idx_glossary_user ON glossary(user_id)")
        cur.execute("CREATE INDEX IF NOT EXISTS idx_history_user ON translation_history(user_id)")

        # 默认管理员账号 admin / admin123
        admin_hash = hash_password("admin123")
        cur.execute(
            "INSERT OR IGNORE INTO users (username, password_hash, is_admin) VALUES (?, ?, 1)",
            ("admin", admin_hash),
        )


# ===========================================================================
# 密码工具
# ===========================================================================

def hash_password(password: str) -> str:
    """SHA-256 哈希密码（加固定盐）。"""
    salt = "ppt_translate_salt_2024"
    return hashlib.sha256((salt + password).encode("utf-8")).hexdigest()


def verify_password(password: str, password_hash: str) -> bool:
    """验证密码。"""
    return hash_password(password) == password_hash


# ===========================================================================
# 认证装饰器
# ===========================================================================

def login_required(f):
    """所有人免登录访问。"""
    @wraps(f)
    def decorated(*args, **kwargs):
        return f(*args, **kwargs)
    return decorated


def get_current_user_id() -> int:
    """获取当前登录用户的 ID；访客模式下返回默认管理员 ID。"""
    uid = session.get("user_id")
    if uid is not None:
        return uid
    # 访客模式：返回 admin 用户 ID
    with db_cursor() as cur:
        cur.execute("SELECT id FROM users WHERE username = ?", ("admin",))
        row = cur.fetchone()
        if row:
            return row[0]
    return 0


# ===========================================================================
# 翻译引擎 TranslationEngine
# ===========================================================================

class TranslationEngine:
    """
    多 AI API 翻译引擎。
    支持 deepseek / doubao / openai 三种接口类型。
    使用 urllib.request 发送 HTTP 请求，不依赖 requests 库。
    """

    DEFAULT_ENDPOINTS = {
        "deepseek": "https://api.deepseek.com/v1/chat/completions",
        "doubao": "https://ark.cn-beijing.volces.com/api/v3/chat/completions",
        "openai": "https://api.openai.com/v1/chat/completions",
    }

    DEFAULT_MODELS = {
        "deepseek": "deepseek-chat",
        "doubao": "doubao-pro-32k",
        "openai": "gpt-3.5-turbo",
    }

    @staticmethod
    def _build_system_prompt(source_lang: str, target_lang: str, glossary_terms: list = None, tm_entries: list = None) -> str:
        """构建翻译系统提示词。"""
        prompt = (
            f"你是一个专业翻译引擎。请将以下 {source_lang} 文本翻译为 {target_lang}。\n"
            f"要求：\n"
            f"1. 保持原文格式和结构\n"
            f"2. 准确传达原文含义\n"
            f"3. 保持术语一致性\n"
            f"4. 只返回翻译结果，不要添加解释或说明\n"
        )

        if glossary_terms:
            prompt += "\n术语表（请务必使用以下术语翻译）：\n"
            for src, tgt in glossary_terms:
                prompt += f"  - {src} -> {tgt}\n"

        if tm_entries:
            prompt += "\n参考翻译记忆：\n"
            for src, tgt in tm_entries:
                prompt += f"  - [{src}] -> [{tgt}]\n"

        return prompt

    @staticmethod
    def _http_post(url: str, headers: dict, body: dict, timeout: int = 120) -> dict:
        """
        使用 urllib.request 发送 POST 请求。
        返回解析后的 JSON 响应。
        """
        import urllib.request
        import urllib.error

        data = json.dumps(body).encode("utf-8")
        req = urllib.request.Request(url, data=data, method="POST")

        for key, value in headers.items():
            req.add_header(key, value)

        try:
            with urllib.request.urlopen(req, timeout=timeout) as resp:
                resp_data = resp.read().decode("utf-8")
                return json.loads(resp_data)
        except urllib.error.HTTPError as e:
            error_body = ""
            try:
                error_body = e.read().decode("utf-8", errors="replace")
            except Exception:
                pass
            logger.error(f"HTTP {e.code}: {error_body}")
            raise Exception(f"API 请求失败 (HTTP {e.code}): {error_body}")
        except urllib.error.URLError as e:
            logger.error(f"URL Error: {e.reason}")
            raise Exception(f"网络请求失败: {e.reason}")
        except Exception as e:
            logger.error(f"请求异常: {str(e)}")
            raise Exception(f"请求异常: {str(e)}")

    @classmethod
    def translate(
        cls,
        source_text: str,
        source_lang: str = "zh",
        target_lang: str = "en",
        api_type: str = "deepseek",
        api_key: str = "",
        model: str = "",
        api_url: str = "",
        glossary_terms: list = None,
        tm_entries: list = None,
    ) -> str:
        """
        调用 AI API 翻译文本。

        参数:
            source_text: 源文本
            source_lang: 源语言代码
            target_lang: 目标语言代码
            api_type: API 类型 (deepseek / doubao / openai)
            api_key: API 密钥
            model: 模型名称
            api_url: 自定义 API 地址（可选）
            glossary_terms: 术语列表 [(source, target), ...]
            tm_entries: 翻译记忆 [(source, target), ...]

        返回:
            翻译后的文本
        """
        if not source_text or not source_text.strip():
            return ""

        # 确定 endpoint
        endpoint = api_url.strip() if api_url else cls.DEFAULT_ENDPOINTS.get(api_type, "")
        if not endpoint:
            raise ValueError(f"未知的 API 类型: {api_type}，且未提供自定义 API 地址")

        # 确定 model
        if not model:
            model = cls.DEFAULT_MODELS.get(api_type, "gpt-3.5-turbo")

        if not api_key:
            raise ValueError("请提供 API Key")

        # 构建请求
        system_prompt = cls._build_system_prompt(
            source_lang, target_lang, glossary_terms, tm_entries
        )

        body = {
            "model": model,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": source_text},
            ],
            "temperature": 0.3,
            "max_tokens": 4096,
        }

        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        }

        logger.info(f"翻译请求: api_type={api_type}, model={model}, text_len={len(source_text)}")

        result = cls._http_post(endpoint, headers, body)

        # 解析响应
        try:
            translated = result["choices"][0]["message"]["content"].strip()
        except (KeyError, IndexError) as e:
            logger.error(f"解析翻译结果失败: {e}, response={result}")
            raise Exception(f"解析翻译结果失败: {result}")

        logger.info(f"翻译完成: result_len={len(translated)}")
        return translated


# ===========================================================================
# 文件解析器 FileParser
# ===========================================================================

class FileParser:
    """
    文件解析器，支持 PPTX 和 PDF 格式。
    返回按页/幻灯片组织的文本结构。
    """

    @staticmethod
    def parse_pptx(file_path: str) -> dict:
        """
        解析 PPTX 文件。
        返回格式:
        {
            "type": "pptx",
            "filename": "xxx.pptx",
            "pages": [
                {"page_number": 1, "paragraphs": ["文本1", "文本2", ...]},
                ...
            ]
        }
        """
        if not HAS_PPTX:
            raise Exception("缺少 python-pptx 库，请运行: pip install python-pptx")

        filename = os.path.basename(file_path)
        prs = Presentation(file_path)
        pages = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            paragraphs = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            paragraphs.append(text)
                # 处理表格
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                paragraphs.append(text)

            if paragraphs:
                pages.append({
                    "page_number": slide_idx,
                    "paragraphs": paragraphs,
                })

        return {
            "type": "pptx",
            "filename": filename,
            "total_slides": len(prs.slides),
            "pages": pages,
        }

    @staticmethod
    def parse_pdf(file_path: str) -> dict:
        """
        解析 PDF 文件。
        返回格式:
        {
            "type": "pdf",
            "filename": "xxx.pdf",
            "pages": [
                {"page_number": 1, "paragraphs": ["文本1", "文本2", ...]},
                ...
            ]
        }
        """
        if not HAS_PDF2:
            raise Exception("缺少 PyPDF2 库，请运行: pip install PyPDF2")

        filename = os.path.basename(file_path)
        reader = PdfReader(file_path)
        pages = []

        for page_idx, page in enumerate(reader.pages, start=1):
            text = page.extract_text()
            if text:
                # 按段落分割
                paragraphs = []
                for line in text.splitlines():
                    line = line.strip()
                    if line:
                        paragraphs.append(line)
                if paragraphs:
                    pages.append({
                        "page_number": page_idx,
                        "paragraphs": paragraphs,
                    })

        return {
            "type": "pdf",
            "filename": filename,
            "total_pages": len(reader.pages),
            "pages": pages,
        }

    @classmethod
    def parse(cls, file_path: str) -> dict:
        """根据文件扩展名自动选择解析器。"""
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pptx":
            return cls.parse_pptx(file_path)
        elif ext == ".pdf":
            return cls.parse_pdf(file_path)
        else:
            raise Exception(f"不支持的文件格式: {ext}，仅支持 .pptx 和 .pdf")


# ===========================================================================
# 翻译记忆库管理 TMManager
# ===========================================================================

class TMManager:
    """翻译记忆库管理器，支持相似度匹配。"""

    @staticmethod
    def _compute_similarity(text_a: str, text_b: str) -> float:
        """
        计算两段文本的相似度（简单字符级 Jaccard + 长度惩罚）。
        返回 0.0 ~ 1.0 的浮点数。
        """
        if not text_a or not text_b:
            return 0.0

        set_a = set(text_a)
        set_b = set(text_b)

        intersection = len(set_a & set_b)
        union = len(set_a | set_b)

        if union == 0:
            return 1.0

        jaccard = intersection / union

        # 长度差异惩罚
        len_ratio = min(len(text_a), len(text_b)) / max(len(text_a), len(text_b))

        return round(jaccard * len_ratio, 4)

    @classmethod
    def find_similar(
        cls,
        source_text: str,
        source_lang: str,
        target_lang: str,
        user_id: int = 0,
        threshold: float = 0.6,
        limit: int = 5,
    ) -> list:
        """
        在翻译记忆库中查找相似翻译。

        返回:
            [(source_text, target_text, similarity), ...]
        """
        with db_cursor() as cur:
            # 查询同语言对的所有 TM 记录
            if user_id > 0:
                cur.execute(
                    "SELECT source_text, target_text FROM translation_memory "
                    "WHERE source_lang=? AND target_lang=? AND (user_id=? OR user_id=0) "
                    "ORDER BY updated_at DESC LIMIT 500",
                    (source_lang, target_lang, user_id),
                )
            else:
                cur.execute(
                    "SELECT source_text, target_text FROM translation_memory "
                    "WHERE source_lang=? AND target_lang=? "
                    "ORDER BY updated_at DESC LIMIT 500",
                    (source_lang, target_lang),
                )

            rows = cur.fetchall()

        results = []
        for row in rows:
            sim = cls._compute_similarity(source_text, row["source_text"])
            if sim >= threshold:
                results.append((row["source_text"], row["target_text"], sim))

        # 按相似度降序排序
        results.sort(key=lambda x: x[2], reverse=True)
        return results[:limit]

    @classmethod
    def add_entry(
        cls,
        user_id: int,
        source_text: str,
        target_text: str,
        source_lang: str,
        target_lang: str,
        similarity: float = 1.0,
    ) -> int:
        """添加翻译记忆条目，返回 ID。"""
        with db_cursor() as cur:
            # 检查是否已存在完全相同的条目
            cur.execute(
                "SELECT id FROM translation_memory "
                "WHERE user_id=? AND source_text=? AND target_text=? AND source_lang=? AND target_lang=?",
                (user_id, source_text, target_text, source_lang, target_lang),
            )
            existing = cur.fetchone()
            if existing:
                cur.execute(
                    "UPDATE translation_memory SET updated_at=datetime('now','localtime'), similarity=? WHERE id=?",
                    (similarity, existing["id"]),
                )
                return existing["id"]

            cur.execute(
                "INSERT INTO translation_memory (user_id, source_text, target_text, source_lang, target_lang, similarity) "
                "VALUES (?, ?, ?, ?, ?, ?)",
                (user_id, source_text, target_text, source_lang, target_lang, similarity),
            )
            return cur.lastrowid


# ===========================================================================
# 语料库管理 GlossaryManager
# ===========================================================================

class GlossaryManager:
    """语料库/术语管理器。"""

    @classmethod
    def get_terms(
        cls,
        user_id: int = 0,
        source_lang: str = "",
        target_lang: str = "",
    ) -> list:
        """获取术语列表。"""
        with db_cursor() as cur:
            if user_id > 0:
                cur.execute(
                    "SELECT * FROM glossary WHERE user_id=? ORDER BY created_at DESC",
                    (user_id,),
                )
            else:
                cur.execute("SELECT * FROM glossary ORDER BY created_at DESC")
            rows = cur.fetchall()
        return [dict(row) for row in rows]

    @classmethod
    def get_terms_for_translation(
        cls,
        source_lang: str,
        target_lang: str,
        user_id: int = 0,
    ) -> list:
        """
        获取翻译用的术语对。
        返回: [(source_term, target_term), ...]
        """
        with db_cursor() as cur:
            if user_id > 0:
                cur.execute(
                    "SELECT source_term, target_term FROM glossary "
                    "WHERE source_lang=? AND target_lang=? AND (user_id=? OR user_id=0)",
                    (source_lang, target_lang, user_id),
                )
            else:
                cur.execute(
                    "SELECT source_term, target_term FROM glossary "
                    "WHERE source_lang=? AND target_lang=?",
                    (source_lang, target_lang),
                )
            rows = cur.fetchall()
        return [(row["source_term"], row["target_term"]) for row in rows]

    @classmethod
    def add_term(
        cls,
        user_id: int,
        source_term: str,
        target_term: str,
        source_lang: str,
        target_lang: str,
        category: str = "",
        note: str = "",
    ) -> int:
        """添加术语，返回 ID。"""
        with db_cursor() as cur:
            cur.execute(
                "INSERT INTO glossary (user_id, source_term, target_term, source_lang, target_lang, category, note) "
                "VALUES (?, ?, ?, ?, ?, ?, ?)",
                (user_id, source_term, target_term, source_lang, target_lang, category, note),
            )
            return cur.lastrowid

    @classmethod
    def delete_term(cls, term_id: int, user_id: int) -> bool:
        """删除术语。"""
        with db_cursor() as cur:
            cur.execute(
                "DELETE FROM glossary WHERE id=? AND (user_id=? OR EXISTS(SELECT 1 FROM users WHERE id=? AND is_admin=1))",
                (term_id, user_id, user_id),
            )
            return cur.rowcount > 0

    @classmethod
    def update_term(
        cls,
        term_id: int,
        user_id: int,
        source_term: str = None,
        target_term: str = None,
        category: str = None,
        note: str = None,
    ) -> bool:
        """更新术语。"""
        fields = []
        params = []
        if source_term is not None:
            fields.append("source_term=?")
            params.append(source_term)
        if target_term is not None:
            fields.append("target_term=?")
            params.append(target_term)
        if category is not None:
            fields.append("category=?")
            params.append(category)
        if note is not None:
            fields.append("note=?")
            params.append(note)

        if not fields:
            return False

        params.extend([term_id, user_id])
        with db_cursor() as cur:
            cur.execute(
                f"UPDATE glossary SET {', '.join(fields)} WHERE id=? AND (user_id=? OR EXISTS(SELECT 1 FROM users WHERE id=? AND is_admin=1))",
                (*params, user_id),
            )
            return cur.rowcount > 0


# ===========================================================================
# 导出器 Exporter
# ===========================================================================

class Exporter:
    """多格式导出器。"""

    @staticmethod
    def export_markdown(data: dict) -> bytes:
        """
        导出为 Markdown 格式。
        data 格式: {"filename": "xxx", "pages": [{"page_number": 1, "paragraphs": [...], "translations": [...]}]}
        """
        lines = []
        lines.append(f"# 翻译结果: {data.get('filename', '未命名')}\n")

        for page in data.get("pages", []):
            page_num = page.get("page_number", "?")
            lines.append(f"## 第 {page_num} 页\n")
            paragraphs = page.get("paragraphs", [])
            translations = page.get("translations", [])

            for i, para in enumerate(paragraphs):
                lines.append(f"**原文 ({i+1}):** {para}\n")
                if i < len(translations):
                    lines.append(f"**译文 ({i+1}):** {translations[i]}\n")
                else:
                    lines.append(f"**译文 ({i+1}):** *(未翻译)*\n")
                lines.append("---\n")

        content = "\n".join(lines)
        return content.encode("utf-8")

    @staticmethod
    def export_json(data: dict) -> bytes:
        """导出为 JSON 格式。"""
        return json.dumps(data, ensure_ascii=False, indent=2).encode("utf-8")

    @staticmethod
    def export_docx(data: dict) -> bytes:
        """导出为 DOCX 格式。"""
        if not HAS_DOCX:
            raise Exception("缺少 python-docx 库，请运行: pip install python-docx")

        doc = Document()
        doc.add_heading(f"翻译结果: {data.get('filename', '未命名')}", level=0)

        for page in data.get("pages", []):
            page_num = page.get("page_number", "?")
            doc.add_heading(f"第 {page_num} 页", level=1)

            paragraphs = page.get("paragraphs", [])
            translations = page.get("translations", [])

            for i, para in enumerate(paragraphs):
                # 原文
                p = doc.add_paragraph()
                run_src = p.add_run(f"原文 ({i+1}): ")
                run_src.bold = True
                p.add_run(para)

                # 译文
                p2 = doc.add_paragraph()
                run_tgt = p2.add_run(f"译文 ({i+1}): ")
                run_tgt.bold = True
                if i < len(translations):
                    run_text = p2.add_run(translations[i])
                    run_text.font.color.rgb = RGBColor(0, 100, 180)
                else:
                    p2.add_run("*(未翻译)*").italic = True

                doc.add_paragraph("---")

        buf = BytesIO()
        doc.save(buf)
        return buf.getvalue()

    @staticmethod
    def export_pdf(data: dict) -> bytes:
        """导出为 PDF 格式。"""
        if not HAS_FPDF:
            raise Exception("缺少 fpdf2 库，请运行: pip install fpdf2")

        pdf = FPDF()
        pdf.add_page()

        # 尝试加载支持中文的字体，如果没有则使用默认字体
        try:
            # 查找系统中的中文字体
            chinese_fonts = [
                "C:/Windows/Fonts/simhei.ttf",    # 黑体
                "C:/Windows/Fonts/simsun.ttc",     # 宋体
                "C:/Windows/Fonts/msyh.ttc",       # 微软雅黑
                "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
                "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
                "/System/Library/Fonts/PingFang.ttc",
            ]
            font_loaded = False
            for font_path in chinese_fonts:
                if os.path.exists(font_path):
                    pdf.add_font("chinese", "", font_path, uni=True)
                    pdf.set_font("chinese", size=12)
                    font_loaded = True
                    break

            if not font_loaded:
                pdf.set_font("Helvetica", size=12)
                logger.warning("未找到中文字体，PDF 可能无法正确显示中文")
        except Exception as e:
            pdf.set_font("Helvetica", size=12)
            logger.warning(f"加载中文字体失败: {e}")

        # 标题
        title = data.get("filename", "Translation Result")
        pdf.set_font_size(18)
        pdf.cell(0, 12, title, new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.ln(6)

        for page in data.get("pages", []):
            page_num = page.get("page_number", "?")
            pdf.set_font_size(14)
            pdf.cell(0, 10, f"Page {page_num}", new_x="LMARGIN", new_y="NEXT")
            pdf.ln(3)

            paragraphs = page.get("paragraphs", [])
            translations = page.get("translations", [])

            pdf.set_font_size(10)
            for i, para in enumerate(paragraphs):
                # 简单处理：先检查当前字体是否支持中文
                pdf.set_font("chinese" if pdf.font_family == "chinese" else "Helvetica", size=10)

                src_text = f"[{i+1}] Source: {para[:200]}"
                pdf.multi_cell(0, 6, src_text)
                pdf.ln(1)

                if i < len(translations):
                    tgt_text = f"[{i+1}] Target: {translations[i][:200]}"
                else:
                    tgt_text = f"[{i+1}] Target: (not translated)"
                pdf.multi_cell(0, 6, tgt_text)
                pdf.ln(3)

                # 检查是否需要换页
                if pdf.get_y() > 260:
                    pdf.add_page()

        return bytes(pdf.output())

    @classmethod
    def export(cls, data: dict, fmt: str) -> tuple:
        """
        导出为指定格式。
        返回 (bytes, content_type, filename)。
        """
        exporters = {
            "md": (cls.export_markdown, "text/markdown; charset=utf-8", "translation.md"),
            "json": (cls.export_json, "application/json; charset=utf-8", "translation.json"),
            "docx": (cls.export_docx, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "translation.docx"),
            "pdf": (cls.export_pdf, "application/pdf", "translation.pdf"),
        }

        if fmt not in exporters:
            raise Exception(f"不支持的导出格式: {fmt}，支持: {', '.join(exporters.keys())}")

        export_func, content_type, filename = exporters[fmt]
        content = export_func(data)
        return content, content_type, filename


# ===========================================================================
# API 配置管理
# ===========================================================================

class ConfigManager:
    """API 配置管理器。"""

    @classmethod
    def get_configs(cls, user_id: int) -> list:
        """获取用户的所有 API 配置。"""
        with db_cursor() as cur:
            cur.execute(
                "SELECT id, api_type, api_key, api_url, model, is_default, created_at "
                "FROM api_configs WHERE user_id=? ORDER BY created_at DESC",
                (user_id,),
            )
            rows = cur.fetchall()
        return [dict(row) for row in rows]

    @classmethod
    def get_default_config(cls, user_id: int, api_type: str = None) -> dict:
        """获取用户的默认 API 配置。"""
        with db_cursor() as cur:
            if api_type:
                cur.execute(
                    "SELECT * FROM api_configs WHERE user_id=? AND api_type=? AND is_default=1 LIMIT 1",
                    (user_id, api_type),
                )
            else:
                cur.execute(
                    "SELECT * FROM api_configs WHERE user_id=? AND is_default=1 LIMIT 1",
                    (user_id,),
                )
            row = cur.fetchone()
        return dict(row) if row else {}

    @classmethod
    def save_config(
        cls,
        user_id: int,
        api_type: str,
        api_key: str = "",
        api_url: str = "",
        model: str = "",
        is_default: int = 0,
    ) -> int:
        """保存 API 配置，返回 ID。"""
        with db_cursor() as cur:
            # 如果设为默认，先取消其他同类型的默认
            if is_default:
                cur.execute(
                    "UPDATE api_configs SET is_default=0 WHERE user_id=? AND api_type=?",
                    (user_id, api_type),
                )

            cur.execute(
                "INSERT INTO api_configs (user_id, api_type, api_key, api_url, model, is_default) "
                "VALUES (?, ?, ?, ?, ?, ?)",
                (user_id, api_type, api_key, api_url, model, is_default),
            )
            return cur.lastrowid

    @classmethod
    def update_config(cls, config_id: int, user_id: int, **kwargs) -> bool:
        """更新 API 配置。"""
        fields = []
        params = []
        for key in ("api_type", "api_key", "api_url", "model", "is_default"):
            if key in kwargs:
                fields.append(f"{key}=?")
                params.append(kwargs[key])

        if not fields:
            return False

        # 如果设为默认，先取消其他同类型的默认
        if kwargs.get("is_default"):
            api_type = kwargs.get("api_type", "")
            if api_type:
                with db_cursor() as cur:
                    cur.execute(
                        "UPDATE api_configs SET is_default=0 WHERE user_id=? AND api_type=? AND id!=?",
                        (user_id, api_type, config_id),
                    )

        params.extend([config_id, user_id])
        with db_cursor() as cur:
            cur.execute(
                f"UPDATE api_configs SET {', '.join(fields)} WHERE id=? AND user_id=?",
                (*params,),
            )
            return cur.rowcount > 0

    @classmethod
    def delete_config(cls, config_id: int, user_id: int) -> bool:
        """删除 API 配置。"""
        with db_cursor() as cur:
            cur.execute(
                "DELETE FROM api_configs WHERE id=? AND user_id=?",
                (config_id, user_id),
            )
            return cur.rowcount > 0


# ===========================================================================
# 翻译历史管理
# ===========================================================================

class HistoryManager:
    """翻译历史管理器。"""

    @classmethod
    def add_record(
        cls,
        user_id: int,
        source_text: str,
        translated_text: str,
        source_lang: str,
        target_lang: str,
        api_type: str = "",
        model: str = "",
        use_tm: int = 1,
    ) -> int:
        """添加翻译历史记录。"""
        with db_cursor() as cur:
            cur.execute(
                "INSERT INTO translation_history (user_id, source_text, translated_text, source_lang, target_lang, api_type, model, use_tm) "
                "VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
                (user_id, source_text, translated_text, source_lang, target_lang, api_type, model, use_tm),
            )
            return cur.lastrowid

    @classmethod
    def get_records(
        cls,
        user_id: int,
        limit: int = 50,
        offset: int = 0,
    ) -> list:
        """获取翻译历史。"""
        with db_cursor() as cur:
            cur.execute(
                "SELECT * FROM translation_history WHERE user_id=? ORDER BY created_at DESC LIMIT ? OFFSET ?",
                (user_id, limit, offset),
            )
            rows = cur.fetchall()
        return [dict(row) for row in rows]

    @classmethod
    def count_records(cls, user_id: int) -> int:
        """获取翻译历史总数。"""
        with db_cursor() as cur:
            cur.execute(
                "SELECT COUNT(*) as cnt FROM translation_history WHERE user_id=?",
                (user_id,),
            )
            row = cur.fetchone()
        return row["cnt"] if row else 0


# ===========================================================================
# 认证相关 API 路由
# ===========================================================================

@app.route("/api/register", methods=["POST"])
def api_register():
    """用户注册。"""
    data = request.get_json(silent=True) or {}
    username = (data.get("username") or "").strip()
    password = (data.get("password") or "").strip()

    if not username or not password:
        return jsonify({"ok": False, "error": "用户名和密码不能为空"})

    if len(username) < 2 or len(username) > 50:
        return jsonify({"ok": False, "error": "用户名长度需在 2-50 个字符之间"})

    if len(password) < 4:
        return jsonify({"ok": False, "error": "密码长度不能少于 4 个字符"})

    try:
        with db_cursor() as cur:
            cur.execute("INSERT INTO users (username, password_hash) VALUES (?, ?)",
                         (username, hash_password(password)))
            user_id = cur.lastrowid
    except sqlite3.IntegrityError:
        return jsonify({"ok": False, "error": "用户名已存在"})

    session["user_id"] = user_id
    session["username"] = username

    logger.info(f"用户注册成功: {username} (ID={user_id})")
    return jsonify({"ok": True, "message": "注册成功", "user_id": user_id, "username": username})


@app.route("/api/login", methods=["POST"])
def api_login():
    """用户登录。"""
    data = request.get_json(silent=True) or {}
    username = (data.get("username") or "").strip()
    password = (data.get("password") or "").strip()

    if not username or not password:
        return jsonify({"ok": False, "error": "用户名和密码不能为空"})

    with db_cursor() as cur:
        cur.execute("SELECT id, username, password_hash, is_admin FROM users WHERE username=?", (username,))
        user = cur.fetchone()

    if not user or not verify_password(password, user["password_hash"]):
        return jsonify({"ok": False, "error": "用户名或密码错误"})

    session["user_id"] = user["id"]
    session["username"] = user["username"]
    session["is_admin"] = bool(user["is_admin"])

    logger.info(f"用户登录成功: {username} (ID={user['id']})")
    return jsonify({
        "ok": True,
        "message": "登录成功",
        "user_id": user["id"],
        "username": user["username"],
        "is_admin": bool(user["is_admin"]),
    })


@app.route("/api/logout", methods=["POST"])
def api_logout():
    """用户登出。"""
    session.clear()
    return jsonify({"ok": True, "message": "已登出"})


@app.route("/api/check_auth", methods=["GET"])
def api_check_auth():
    """检查认证状态。"""
    if "user_id" in session:
        return jsonify({
            "ok": True,
            "authenticated": True,
            "user_id": session["user_id"],
            "username": session.get("username", ""),
            "is_admin": session.get("is_admin", False),
        })
    return jsonify({"ok": True, "authenticated": False})


# ===========================================================================
# 文件上传与解析 API
# ===========================================================================

@app.route("/api/upload", methods=["POST"])
@login_required
def api_upload():
    """
    上传 PPTX/PDF 文件并解析。
    表单字段: file (文件), source_lang (可选), target_lang (可选)
    返回解析后的段落结构。
    """
    if "file" not in request.files:
        return jsonify({"ok": False, "error": "未提供文件"})

    file = request.files["file"]
    if not file.filename:
        return jsonify({"ok": False, "error": "文件名为空"})

    # 验证文件类型
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in (".pptx", ".pdf"):
        return jsonify({"ok": False, "error": "仅支持 .pptx 和 .pdf 文件"})

    # 检查依赖
    if ext == ".pptx" and not HAS_PPTX:
        return jsonify({"ok": False, "error": "服务器缺少 python-pptx 库，请联系管理员安装"})
    if ext == ".pdf" and not HAS_PDF2:
        return jsonify({"ok": False, "error": "服务器缺少 PyPDF2 库，请联系管理员安装"})

    # 保存文件
    unique_name = f"{uuid.uuid4().hex[:12]}_{file.filename}"
    file_path = os.path.join(UPLOAD_FOLDER, unique_name)
    file.save(file_path)

    try:
        result = FileParser.parse(file_path)
        result["file_id"] = unique_name
        result["source_lang"] = request.form.get("source_lang", "zh")
        result["target_lang"] = request.form.get("target_lang", "en")

        # 为每个段落生成唯一 ID
        para_idx = 0
        for page in result["pages"]:
            new_paragraphs = []
            for para in page["paragraphs"]:
                new_paragraphs.append({
                    "id": f"para_{para_idx}",
                    "text": para,
                    "translation": "",
                })
                para_idx += 1
            page["paragraphs"] = new_paragraphs

        logger.info(f"文件解析成功: {file.filename}, 段落数={para_idx}")
        return jsonify({"ok": True, "data": result})

    except Exception as e:
        logger.error(f"文件解析失败: {e}\n{traceback.format_exc()}")
        return jsonify({"ok": False, "error": f"文件解析失败: {str(e)}"})
    finally:
        # 不删除文件，后续翻译/导出可能需要
        pass


# ===========================================================================
# 翻译 API
# ===========================================================================

@app.route("/api/translate", methods=["POST"])
@login_required
def api_translate():
    """
    翻译单条文本。

    JSON Body:
        source_text: 源文本
        source_lang: 源语言 (默认 "zh")
        target_lang: 目标语言 (默认 "en")
        api_type: API 类型 (deepseek / doubao / openai)
        api_key: API 密钥
        model: 模型名称 (可选)
        api_url: 自定义 API 地址 (可选)
        use_tm: 是否使用翻译记忆 (默认 true)
        save_tm: 是否保存到翻译记忆 (默认 true)
    """
    data = request.get_json(silent=True) or {}
    source_text = (data.get("source_text") or "").strip()
    source_lang = data.get("source_lang", "zh").strip()
    target_lang = data.get("target_lang", "en").strip()
    api_type = data.get("api_type", "deepseek").strip()
    api_key = data.get("api_key", "").strip()
    model = data.get("model", "").strip()
    api_url = data.get("api_url", "").strip()
    use_tm = data.get("use_tm", True)
    save_tm = data.get("save_tm", True)

    if not source_text:
        return jsonify({"ok": False, "error": "源文本不能为空"})

    user_id = get_current_user_id()

    # 如果未提供 api_key，尝试从数据库获取默认配置
    if not api_key:
        config = ConfigManager.get_default_config(user_id, api_type)
        if config:
            api_key = config.get("api_key", "")
            model = model or config.get("model", "")
            api_url = api_url or config.get("api_url", "")

    try:
        # 获取术语和翻译记忆
        glossary_terms = GlossaryManager.get_terms_for_translation(source_lang, target_lang, user_id)
        tm_entries = []
        if use_tm:
            tm_results = TMManager.find_similar(source_text, source_lang, target_lang, user_id)
            tm_entries = [(src, tgt) for src, tgt, sim in tm_results]

        # 调用翻译引擎
        translated = TranslationEngine.translate(
            source_text=source_text,
            source_lang=source_lang,
            target_lang=target_lang,
            api_type=api_type,
            api_key=api_key,
            model=model,
            api_url=api_url,
            glossary_terms=glossary_terms,
            tm_entries=tm_entries,
        )

        # 保存翻译记忆
        if save_tm and translated:
            TMManager.add_entry(
                user_id=user_id,
                source_text=source_text,
                target_text=translated,
                source_lang=source_lang,
                target_lang=target_lang,
                similarity=1.0,
            )

        # 保存翻译历史
        HistoryManager.add_record(
            user_id=user_id,
            source_text=source_text,
            translated_text=translated,
            source_lang=source_lang,
            target_lang=target_lang,
            api_type=api_type,
            model=model,
            use_tm=1 if use_tm else 0,
        )

        return jsonify({
            "ok": True,
            "data": {
                "source_text": source_text,
                "translated_text": translated,
                "source_lang": source_lang,
                "target_lang": target_lang,
                "api_type": api_type,
                "model": model,
                "tm_used": len(tm_entries) > 0,
                "glossary_used": len(glossary_terms) > 0,
            },
        })

    except Exception as e:
        logger.error(f"翻译失败: {e}\n{traceback.format_exc()}")
        return jsonify({"ok": False, "error": f"翻译失败: {str(e)}"})


@app.route("/api/translate/batch", methods=["POST"])
@login_required
def api_translate_batch():
    """
    批量翻译段落。

    JSON Body:
        paragraphs: [{"id": "para_0", "text": "..."}, ...]
        source_lang: 源语言
        target_lang: 目标语言
        api_type: API 类型
        api_key: API 密钥
        model: 模型名称 (可选)
        api_url: 自定义 API 地址 (可选)
        use_tm: 是否使用翻译记忆 (默认 true)
    """
    data = request.get_json(silent=True) or {}
    paragraphs = data.get("paragraphs", [])
    source_lang = data.get("source_lang", "zh").strip()
    target_lang = data.get("target_lang", "en").strip()
    api_type = data.get("api_type", "deepseek").strip()
    api_key = data.get("api_key", "").strip()
    model = data.get("model", "").strip()
    api_url = data.get("api_url", "").strip()
    use_tm = data.get("use_tm", True)

    if not paragraphs:
        return jsonify({"ok": False, "error": "没有提供要翻译的段落"})

    if not isinstance(paragraphs, list):
        return jsonify({"ok": False, "error": "paragraphs 必须是数组"})

    user_id = get_current_user_id()

    if not api_key:
        config = ConfigManager.get_default_config(user_id, api_type)
        if config:
            api_key = config.get("api_key", "")
            model = model or config.get("model", "")
            api_url = api_url or config.get("api_url", "")

    results = []
    errors = []

    # 获取术语
    glossary_terms = GlossaryManager.get_terms_for_translation(source_lang, target_lang, user_id)

    for para in paragraphs:
        para_id = para.get("id", "")
        text = (para.get("text") or "").strip()
        if not text:
            results.append({"id": para_id, "text": text, "translation": "", "status": "skipped"})
            continue

        try:
            # TM 匹配
            tm_entries = []
            if use_tm:
                tm_results = TMManager.find_similar(text, source_lang, target_lang, user_id)
                tm_entries = [(src, tgt) for src, tgt, sim in tm_results]
                # 如果有完全匹配，直接使用
                for src, tgt, sim in tm_results:
                    if sim >= 0.95:
                        results.append({
                            "id": para_id,
                            "text": text,
                            "translation": tgt,
                            "status": "tm_match",
                        })
                        # 保存历史
                        HistoryManager.add_record(
                            user_id=user_id, source_text=text, translated_text=tgt,
                            source_lang=source_lang, target_lang=target_lang,
                            api_type="tm", model="", use_tm=1,
                        )
                        break
                else:
                    # 无高相似度匹配，调用 API
                    translated = TranslationEngine.translate(
                        source_text=text, source_lang=source_lang, target_lang=target_lang,
                        api_type=api_type, api_key=api_key, model=model, api_url=api_url,
                        glossary_terms=glossary_terms, tm_entries=tm_entries,
                    )
                    results.append({
                        "id": para_id,
                        "text": text,
                        "translation": translated,
                        "status": "translated",
                    })
                    # 保存 TM
                    TMManager.add_entry(user_id, text, translated, source_lang, target_lang)
                    HistoryManager.add_record(
                        user_id=user_id, source_text=text, translated_text=translated,
                        source_lang=source_lang, target_lang=target_lang,
                        api_type=api_type, model=model, use_tm=1 if use_tm else 0,
                    )
                    continue
                continue
        except Exception as e:
            error_msg = str(e)
            errors.append({"id": para_id, "error": error_msg})
            results.append({"id": para_id, "text": text, "translation": "", "status": "error", "error": error_msg})
            continue

    success_count = sum(1 for r in results if r["status"] in ("translated", "tm_match"))
    logger.info(f"批量翻译完成: 总数={len(paragraphs)}, 成功={success_count}, 失败={len(errors)}")

    return jsonify({
        "ok": True,
        "data": {
            "results": results,
            "total": len(paragraphs),
            "success": success_count,
            "errors": len(errors),
            "error_details": errors,
        },
    })


@app.route("/api/translate/file", methods=["POST"])
@login_required
def api_translate_file():
    """
    一站式文件翻译：上传 + 自动翻译整个文件。

    表单字段:
        file: PPTX/PDF 文件
        source_lang: 源语言 (默认 "zh")
        target_lang: 目标语言 (默认 "en")
        api_type: API 类型 (默认 "deepseek")
        api_key: API 密钥
        model: 模型名称 (可选)
        api_url: 自定义 API 地址 (可选)
        use_tm: 是否使用翻译记忆 (默认 true)
    """
    if "file" not in request.files:
        return jsonify({"ok": False, "error": "未提供文件"})

    file = request.files["file"]
    if not file.filename:
        return jsonify({"ok": False, "error": "文件名为空"})

    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in (".pptx", ".pdf"):
        return jsonify({"ok": False, "error": "仅支持 .pptx 和 .pdf 文件"})

    if ext == ".pptx" and not HAS_PPTX:
        return jsonify({"ok": False, "error": "服务器缺少 python-pptx 库"})
    if ext == ".pdf" and not HAS_PDF2:
        return jsonify({"ok": False, "error": "服务器缺少 PyPDF2 库"})

    source_lang = request.form.get("source_lang", "zh").strip()
    target_lang = request.form.get("target_lang", "en").strip()
    api_type = request.form.get("api_type", "deepseek").strip()
    api_key = request.form.get("api_key", "").strip()
    model = request.form.get("model", "").strip()
    api_url = request.form.get("api_url", "").strip()
    use_tm = request.form.get("use_tm", "true").lower() in ("true", "1", "yes")

    # 保存文件
    unique_name = f"{uuid.uuid4().hex[:12]}_{file.filename}"
    file_path = os.path.join(UPLOAD_FOLDER, unique_name)
    file.save(file_path)

    user_id = get_current_user_id()

    if not api_key:
        config = ConfigManager.get_default_config(user_id, api_type)
        if config:
            api_key = config.get("api_key", "")
            model = model or config.get("model", "")
            api_url = api_url or config.get("api_url", "")

    try:
        # 解析文件
        parsed = FileParser.parse(file_path)

        # 收集所有段落
        all_paragraphs = []
        for page in parsed["pages"]:
            for para in page["paragraphs"]:
                all_paragraphs.append(para)

        # 获取术语
        glossary_terms = GlossaryManager.get_terms_for_translation(source_lang, target_lang, user_id)

        # 逐段翻译
        translated_pages = []
        global_para_idx = 0
        for page in parsed["pages"]:
            page_translations = []
            for para in page["paragraphs"]:
                try:
                    # TM 匹配
                    tm_entries = []
                    if use_tm:
                        tm_results = TMManager.find_similar(para, source_lang, target_lang, user_id)
                        tm_entries = [(src, tgt) for src, tgt, sim in tm_results]
                        for src, tgt, sim in tm_results:
                            if sim >= 0.95:
                                page_translations.append(tgt)
                                HistoryManager.add_record(
                                    user_id=user_id, source_text=para, translated_text=tgt,
                                    source_lang=source_lang, target_lang=target_lang,
                                    api_type="tm", model="", use_tm=1,
                                )
                                break
                        else:
                            translated = TranslationEngine.translate(
                                source_text=para, source_lang=source_lang, target_lang=target_lang,
                                api_type=api_type, api_key=api_key, model=model, api_url=api_url,
                                glossary_terms=glossary_terms, tm_entries=tm_entries,
                            )
                            page_translations.append(translated)
                            TMManager.add_entry(user_id, para, translated, source_lang, target_lang)
                            HistoryManager.add_record(
                                user_id=user_id, source_text=para, translated_text=translated,
                                source_lang=source_lang, target_lang=target_lang,
                                api_type=api_type, model=model, use_tm=1 if use_tm else 0,
                            )
                            continue
                        continue
                except Exception as e:
                    logger.error(f"翻译段落 {global_para_idx} 失败: {e}")
                    page_translations.append(f"[翻译失败: {str(e)}]")

                global_para_idx += 1

            translated_pages.append({
                "page_number": page["page_number"],
                "paragraphs": page["paragraphs"],
                "translations": page_translations,
            })

        result = {
            "file_id": unique_name,
            "type": parsed["type"],
            "filename": parsed["filename"],
            "source_lang": source_lang,
            "target_lang": target_lang,
            "total_pages": parsed.get("total_pages") or parsed.get("total_slides", 0),
            "pages": translated_pages,
        }

        logger.info(f"文件翻译完成: {file.filename}, 总段落数={global_para_idx}")
        return jsonify({"ok": True, "data": result})

    except Exception as e:
        logger.error(f"文件翻译失败: {e}\n{traceback.format_exc()}")
        return jsonify({"ok": False, "error": f"文件翻译失败: {str(e)}"})


# ===========================================================================
# 语料库/术语 API
# ===========================================================================

@app.route("/api/glossary", methods=["GET"])
@login_required
def api_glossary_get():
    """获取术语列表。"""
    user_id = get_current_user_id()
    source_lang = request.args.get("source_lang", "").strip()
    target_lang = request.args.get("target_lang", "").strip()
    terms = GlossaryManager.get_terms(user_id, source_lang, target_lang)
    return jsonify({"ok": True, "data": terms})


@app.route("/api/glossary", methods=["POST"])
@login_required
def api_glossary_post():
    """添加术语。"""
    data = request.get_json(silent=True) or {}
    user_id = get_current_user_id()

    source_term = (data.get("source_term") or "").strip()
    target_term = (data.get("target_term") or "").strip()
    source_lang = (data.get("source_lang") or "").strip()
    target_lang = (data.get("target_lang") or "").strip()
    category = (data.get("category") or "").strip()
    note = (data.get("note") or "").strip()

    if not source_term or not target_term:
        return jsonify({"ok": False, "error": "源术语和目标术语不能为空"})

    if not source_lang or not target_lang:
        return jsonify({"ok": False, "error": "语言对不能为空"})

    term_id = GlossaryManager.add_term(
        user_id=user_id,
        source_term=source_term,
        target_term=target_term,
        source_lang=source_lang,
        target_lang=target_lang,
        category=category,
        note=note,
    )

    return jsonify({"ok": True, "message": "术语添加成功", "term_id": term_id})


@app.route("/api/glossary", methods=["DELETE"])
@login_required
def api_glossary_delete():
    """删除术语。"""
    data = request.get_json(silent=True) or {}
    term_id = data.get("term_id")
    if not term_id:
        return jsonify({"ok": False, "error": "缺少 term_id"})

    user_id = get_current_user_id()
    ok = GlossaryManager.delete_term(term_id, user_id)
    if ok:
        return jsonify({"ok": True, "message": "术语已删除"})
    return jsonify({"ok": False, "error": "删除失败，术语不存在或无权限"})


@app.route("/api/glossary/update", methods=["POST"])
@login_required
def api_glossary_update():
    """更新术语。"""
    data = request.get_json(silent=True) or {}
    term_id = data.get("term_id")
    if not term_id:
        return jsonify({"ok": False, "error": "缺少 term_id"})

    user_id = get_current_user_id()
    ok = GlossaryManager.update_term(
        term_id=term_id,
        user_id=user_id,
        source_term=data.get("source_term"),
        target_term=data.get("target_term"),
        category=data.get("category"),
        note=data.get("note"),
    )
    if ok:
        return jsonify({"ok": True, "message": "术语已更新"})
    return jsonify({"ok": False, "error": "更新失败，术语不存在或无权限"})


# ===========================================================================
# 翻译记忆库 API
# ===========================================================================

@app.route("/api/tm", methods=["GET"])
@login_required
def api_tm_get():
    """获取翻译记忆列表。"""
    user_id = get_current_user_id()
    source_lang = request.args.get("source_lang", "").strip()
    target_lang = request.args.get("target_lang", "").strip()
    search_text = request.args.get("search", "").strip()
    limit = request.args.get("limit", 100, type=int)

    with db_cursor() as cur:
        query = "SELECT * FROM translation_memory WHERE (user_id=? OR user_id=0)"
        params = [user_id]

        if source_lang and target_lang:
            query += " AND source_lang=? AND target_lang=?"
            params.extend([source_lang, target_lang])

        if search_text:
            query += " AND (source_text LIKE ? OR target_text LIKE ?)"
            params.extend([f"%{search_text}%", f"%{search_text}%"])

        query += " ORDER BY updated_at DESC LIMIT ?"
        params.append(limit)

        cur.execute(query, params)
        rows = cur.fetchall()

    return jsonify({"ok": True, "data": [dict(row) for row in rows]})


@app.route("/api/tm", methods=["POST"])
@login_required
def api_tm_post():
    """添加翻译记忆条目。"""
    data = request.get_json(silent=True) or {}
    user_id = get_current_user_id()

    source_text = (data.get("source_text") or "").strip()
    target_text = (data.get("target_text") or "").strip()
    source_lang = (data.get("source_lang") or "").strip()
    target_lang = (data.get("target_lang") or "").strip()
    similarity = float(data.get("similarity", 1.0))

    if not source_text or not target_text:
        return jsonify({"ok": False, "error": "源文本和目标文本不能为空"})

    tm_id = TMManager.add_entry(
        user_id=user_id,
        source_text=source_text,
        target_text=target_text,
        source_lang=source_lang,
        target_lang=target_lang,
        similarity=similarity,
    )

    return jsonify({"ok": True, "message": "翻译记忆已添加", "tm_id": tm_id})


@app.route("/api/tm/search", methods=["POST"])
@login_required
def api_tm_search():
    """搜索相似翻译记忆。"""
    data = request.get_json(silent=True) or {}
    user_id = get_current_user_id()

    source_text = (data.get("source_text") or "").strip()
    source_lang = (data.get("source_lang") or "zh").strip()
    target_lang = (data.get("target_lang") or "en").strip()
    threshold = float(data.get("threshold", 0.6))
    limit = int(data.get("limit", 5))

    if not source_text:
        return jsonify({"ok": False, "error": "搜索文本不能为空"})

    results = TMManager.find_similar(
        source_text=source_text,
        source_lang=source_lang,
        target_lang=target_lang,
        user_id=user_id,
        threshold=threshold,
        limit=limit,
    )

    return jsonify({
        "ok": True,
        "data": [
            {"source_text": src, "target_text": tgt, "similarity": sim}
            for src, tgt, sim in results
        ],
    })


@app.route("/api/tm", methods=["DELETE"])
@login_required
def api_tm_delete():
    """删除翻译记忆条目。"""
    data = request.get_json(silent=True) or {}
    tm_id = data.get("tm_id")
    if not tm_id:
        return jsonify({"ok": False, "error": "缺少 tm_id"})

    user_id = get_current_user_id()
    with db_cursor() as cur:
        cur.execute(
            "DELETE FROM translation_memory WHERE id=? AND (user_id=? OR EXISTS(SELECT 1 FROM users WHERE id=? AND is_admin=1))",
            (tm_id, user_id, user_id),
        )
        ok = cur.rowcount > 0

    if ok:
        return jsonify({"ok": True, "message": "翻译记忆已删除"})
    return jsonify({"ok": False, "error": "删除失败，记录不存在或无权限"})


@app.route("/api/tm/clear", methods=["POST"])
@login_required
def api_tm_clear():
    """清空当前用户的翻译记忆。"""
    user_id = get_current_user_id()
    with db_cursor() as cur:
        cur.execute("DELETE FROM translation_memory WHERE user_id=?", (user_id,))
    return jsonify({"ok": True, "message": "翻译记忆已清空"})


# ===========================================================================
# 导出 API
# ===========================================================================

@app.route("/api/export/<fmt>", methods=["GET"])
@login_required
def api_export(fmt: str):
    """
    导出翻译结果为指定格式。

    查询参数:
        file_id: 文件 ID（上传时返回的 unique_name）
        或者通过 JSON body 传递完整数据（POST 请求）

    支持格式: md, json, docx, pdf
    """
    file_id = request.args.get("file_id", "").strip()

    if not file_id:
        # 尝试从 POST body 获取
        data = request.get_json(silent=True) or {}
        if data.get("data"):
            export_data = data["data"]
        else:
            return jsonify({"ok": False, "error": "请提供 file_id 或导出数据"})
    else:
        # 检查是否是 POST 请求带有覆盖数据
        data = request.get_json(silent=True) or {}
        export_data = data.get("data", None)

        if not export_data:
            # 返回示例/空模板（用户应在前端传递完整数据）
            return jsonify({"ok": False, "error": "请在 POST body 中提供 data 字段（包含翻译结果）"})

    try:
        content, content_type, filename = Exporter.export(export_data, fmt)

        response = make_response(content)
        response.headers["Content-Type"] = content_type
        response.headers["Content-Disposition"] = f"attachment; filename={filename}"
        response.headers["Content-Length"] = len(content)
        return response

    except Exception as e:
        logger.error(f"导出失败: {e}\n{traceback.format_exc()}")
        return jsonify({"ok": False, "error": f"导出失败: {str(e)}"})


@app.route("/api/export/<fmt>", methods=["POST"])
@login_required
def api_export_post(fmt: str):
    """POST 方式导出翻译结果。"""
    data = request.get_json(silent=True) or {}
    export_data = data.get("data")

    if not export_data:
        return jsonify({"ok": False, "error": "请在 body 中提供 data 字段"})

    try:
        content, content_type, filename = Exporter.export(export_data, fmt)

        response = make_response(content)
        response.headers["Content-Type"] = content_type
        response.headers["Content-Disposition"] = f"attachment; filename={filename}"
        response.headers["Content-Length"] = len(content)
        return response

    except Exception as e:
        logger.error(f"导出失败: {e}\n{traceback.format_exc()}")
        return jsonify({"ok": False, "error": f"导出失败: {str(e)}"})


# ===========================================================================
# API 配置管理 API
# ===========================================================================

@app.route("/api/config", methods=["GET"])
@login_required
def api_config_get():
    """获取当前用户的 API 配置列表。"""
    user_id = get_current_user_id()
    configs = ConfigManager.get_configs(user_id)
    # 隐藏 API Key 的部分内容
    for cfg in configs:
        key = cfg.get("api_key", "")
        if key and len(key) > 8:
            cfg["api_key_masked"] = key[:4] + "****" + key[-4:]
        else:
            cfg["api_key_masked"] = "****" if key else ""
    return jsonify({"ok": True, "data": configs})


@app.route("/api/config", methods=["POST"])
@login_required
def api_config_post():
    """保存 API 配置。"""
    data = request.get_json(silent=True) or {}
    user_id = get_current_user_id()

    api_type = (data.get("api_type") or "").strip()
    api_key = (data.get("api_key") or "").strip()
    api_url = (data.get("api_url") or "").strip()
    model = (data.get("model") or "").strip()
    is_default = 1 if data.get("is_default") else 0

    if not api_type:
        return jsonify({"ok": False, "error": "API 类型不能为空"})

    if api_type not in TranslationEngine.DEFAULT_ENDPOINTS and not api_url:
        return jsonify({"ok": False, "error": f"不支持的 API 类型: {api_type}，需提供自定义 API 地址"})

    config_id = ConfigManager.save_config(
        user_id=user_id,
        api_type=api_type,
        api_key=api_key,
        api_url=api_url,
        model=model,
        is_default=is_default,
    )

    return jsonify({"ok": True, "message": "配置已保存", "config_id": config_id})


@app.route("/api/config", methods=["DELETE"])
@login_required
def api_config_delete():
    """删除 API 配置。"""
    data = request.get_json(silent=True) or {}
    config_id = data.get("config_id")
    if not config_id:
        return jsonify({"ok": False, "error": "缺少 config_id"})

    user_id = get_current_user_id()
    ok = ConfigManager.delete_config(config_id, user_id)
    if ok:
        return jsonify({"ok": True, "message": "配置已删除"})
    return jsonify({"ok": False, "error": "删除失败"})


@app.route("/api/config", methods=["PUT"])
@login_required
def api_config_put():
    """更新 API 配置。"""
    data = request.get_json(silent=True) or {}
    config_id = data.get("config_id")
    if not config_id:
        return jsonify({"ok": False, "error": "缺少 config_id"})

    user_id = get_current_user_id()
    ok = ConfigManager.update_config(
        config_id=config_id,
        user_id=user_id,
        api_type=data.get("api_type"),
        api_key=data.get("api_key"),
        api_url=data.get("api_url"),
        model=data.get("model"),
        is_default=1 if data.get("is_default") else 0,
    )
    if ok:
        return jsonify({"ok": True, "message": "配置已更新"})
    return jsonify({"ok": False, "error": "更新失败"})


# ===========================================================================
# 翻译历史 API
# ===========================================================================

@app.route("/api/history", methods=["GET"])
@login_required
def api_history_get():
    """获取翻译历史。"""
    user_id = get_current_user_id()
    limit = request.args.get("limit", 50, type=int)
    offset = request.args.get("offset", 0, type=int)

    records = HistoryManager.get_records(user_id, limit, offset)
    total = HistoryManager.count_records(user_id)

    return jsonify({
        "ok": True,
        "data": records,
        "total": total,
        "limit": limit,
        "offset": offset,
    })


@app.route("/api/history/clear", methods=["POST"])
@login_required
def api_history_clear():
    """清空翻译历史。"""
    user_id = get_current_user_id()
    with db_cursor() as cur:
        cur.execute("DELETE FROM translation_history WHERE user_id=?", (user_id,))
    return jsonify({"ok": True, "message": "翻译历史已清空"})


# ===========================================================================
# 依赖状态 API
# ===========================================================================

@app.route("/api/status", methods=["GET"])
def api_status():
    """获取服务器状态和依赖信息。"""
    deps = {
        "python_pptx": HAS_PPTX,
        "pypdf2": HAS_PDF2,
        "python_docx": HAS_DOCX,
        "fpdf2": HAS_FPDF,
    }

    upload_dir_ok = os.path.isdir(UPLOAD_FOLDER)
    db_ok = os.path.isfile(DATABASE_PATH)

    return jsonify({
        "ok": True,
        "status": "running",
        "version": "1.0.0",
        "dependencies": deps,
        "upload_dir_ok": upload_dir_ok,
        "database_ok": db_ok,
        "supported_api_types": list(TranslationEngine.DEFAULT_ENDPOINTS.keys()),
        "supported_export_formats": ["md", "json", "docx", "pdf"],
    })


# ===========================================================================
# 首页（简单提示页面）
# ===========================================================================

@app.route("/", methods=["GET"])
@app.route("/workspace", methods=["GET"])
def index():
    """返回前端页面。"""
    return render_template("index.html")


# ===========================================================================
# 错误处理
# ===========================================================================

@app.errorhandler(404)
def handle_404(e):
    return jsonify({"ok": False, "error": "接口不存在"}), 404


@app.errorhandler(405)
def handle_405(e):
    return jsonify({"ok": False, "error": "请求方法不允许"}), 405


@app.errorhandler(413)
def handle_413(e):
    return jsonify({"ok": False, "error": "文件大小超出限制（最大 100MB）"}), 413


@app.errorhandler(500)
def handle_500(e):
    logger.error(f"服务器错误: {e}\n{traceback.format_exc()}")
    return jsonify({"ok": False, "error": "服务器内部错误"}), 500


# 模块级别初始化数据库（gunicorn 导入时执行）
init_db()


# ===========================================================================
# 启动入口
# ===========================================================================

if __name__ == "__main__":
    logger.info("=" * 60)
    logger.info("PPT/PDF 翻译工作台 启动中...")
    logger.info("=" * 60)

    # 打印依赖状态
    logger.info(f"python-pptx: {'已安装' if HAS_PPTX else '未安装 (pip install python-pptx)'}")
    logger.info(f"PyPDF2:     {'已安装' if HAS_PDF2 else '未安装 (pip install PyPDF2)'}")
    logger.info(f"python-docx: {'已安装' if HAS_DOCX else '未安装 (pip install python-docx)'}")
    logger.info(f"fpdf2:      {'已安装' if HAS_FPDF else '未安装 (pip install fpdf2)'}")

    # 初始化数据库
    init_db()
    logger.info(f"数据库初始化完成: {DATABASE_PATH}")
    logger.info(f"上传目录: {UPLOAD_FOLDER}")

    logger.info("-" * 60)
    logger.info(f"监听地址: http://0.0.0.0:5001")
    logger.info(f"默认管理员: admin / admin123")
    logger.info("-" * 60)

    app.run(host="0.0.0.0", port=5001, debug=False)
